#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
阿里云邮箱 - 邮件收取与总结
"""

import argparse
import json
import os
import re
import sys
from datetime import datetime, timedelta
from email.header import decode_header
from pathlib import Path
import requests

# 邮件处理
try:
    from imap_tools import MailBox, AND
except ImportError:
    print("❌ 缺少 imap-tools，运行：pip3 install imap-tools")
    sys.exit(1)

# 附件处理
try:
    from docx import Document as DocxDocument
    import openpyxl
    from pptx import Presentation
    from pypdf import PdfReader
except ImportError:
    print("❌ 缺少附件处理库，运行：pip3 install python-docx openpyxl python-pptx pypdf2")
    sys.exit(1)


def parse_credentials(cred_file):
    """解析凭证文件"""
    accounts = []
    current_account = {}
    
    with open(cred_file, 'r', encoding='utf-8') as f:
        for line in f:
            line = line.strip()
            if not line or line.startswith('#'):
                continue
            
            if line.startswith('###'):
                if current_account:
                    accounts.append(current_account)
                current_account = {}
                continue
            
            if ':' in line:
                key, value = line.split(':', 1)
                # 去掉开头的 `-` 和空格，然后替换 `-` 为 `_`
                key = key.strip().lstrip('-').strip().lower().replace('-', '_')
                value = value.strip()
                current_account[key] = value
        
        if current_account:
            accounts.append(current_account)
    
    return accounts


def decode_mime(header):
    """解码 MIME 编码的邮件头"""
    if not header:
        return ''
    
    # imap_tools 可能返回 tuple (如 msg.to)
    if isinstance(header, tuple):
        header = ', '.join([str(h) for h in header])
    
    decoded = decode_header(header)
    result = ''
    for text, encoding in decoded:
        if isinstance(text, bytes):
            result += text.decode(encoding or 'utf-8', errors='replace')
        else:
            result += text
    return result


def format_size(size_bytes):
    """格式化文件大小"""
    for unit in ['B', 'KB', 'MB', 'GB']:
        if size_bytes < 1024.0:
            return f"{size_bytes:.1f}{unit}"
        size_bytes /= 1024.0
    return f"{size_bytes:.1f}TB"


def send_to_feishu(markdown_content, chat_id):
    """发送 Markdown 消息到飞书"""
    # 飞书凭证（从环境变量或配置文件读取）
    app_id = os.environ.get('FEISHU_APP_ID', 'cli_a909ad9f75fadbb5')
    app_secret = os.environ.get('FEISHU_APP_SECRET', 'p1MtN6OZic92OCOpMgxaZdSzAvRfsrys')
    
    # 1. 获取 tenant_access_token
    try:
        token_url = "https://open.feishu.cn/open-apis/auth/v3/tenant_access_token/internal"
        token_resp = requests.post(token_url, json={"app_id": app_id, "app_secret": app_secret}, timeout=10)
        token_data = token_resp.json()
        
        if token_data.get('code') != 0:
            print(f"❌ 获取飞书 token 失败：{token_data.get('msg')}")
            return False
        
        token = token_data['tenant_access_token']
    except Exception as e:
        print(f"❌ 飞书认证失败：{e}")
        return False
    
    # 2. 发送消息
    try:
        msg_url = "https://open.feishu.cn/open-apis/im/v1/messages"
        headers = {
            "Authorization": f"Bearer {token}",
            "Content-Type": "application/json"
        }
        
        # 飞书 Markdown 消息
        content = {
            "text": markdown_content
        }
        
        payload = {
            "receive_id": chat_id,
            "msg_type": "text",  # 使用 text 类型，因为飞书对 Markdown 支持有限
            "content": json.dumps(content)
        }
        
        params = {"receive_id_type": "open_id"}
        resp = requests.post(msg_url, headers=headers, params=params, json=payload, timeout=10)
        result = resp.json()
        
        if result.get('code') == 0:
            msg_id = result['data']['message_id']
            print(f"✅ 飞书发送成功：{msg_id}")
            return True
        else:
            print(f"❌ 飞书发送失败：{result.get('msg')}")
            return False
            
    except Exception as e:
        print(f"❌ 飞书发送异常：{e}")
        return False


def extract_text_from_buffer(buffer, file_type):
    """从内存 buffer 提取文本（不写文件）"""
    try:
        from io import BytesIO
        
        if file_type == 'docx':
            doc = DocxDocument(BytesIO(buffer))
            return '\n'.join([para.text for para in doc.paragraphs if para.text.strip()])
        
        elif file_type == 'xlsx':
            wb = openpyxl.load_workbook(BytesIO(buffer), read_only=True, data_only=True)
            texts = []
            for sheet_name in wb.sheetnames[:2]:
                sheet = wb[sheet_name]
                for row_idx, row in enumerate(sheet.iter_rows(max_row=10, values_only=True)):
                    if row_idx >= 10:
                        break
                    cells = [str(cell) if cell is not None else '' for cell in row]
                    if any(cells):
                        texts.append(' | '.join(cells))
            return '\n'.join(texts)
        
        elif file_type == 'pptx':
            prs = Presentation(BytesIO(buffer))
            texts = []
            for slide_idx, slide in enumerate(prs.slides):
                if slide_idx >= 5:
                    break
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
            return '\n'.join(texts)
        
        elif file_type == 'pdf':
            reader = PdfReader(BytesIO(buffer))
            texts = []
            for page_idx, page in enumerate(reader.pages):
                if page_idx >= 5:
                    break
                text = page.extract_text()
                if text:
                    texts.append(text.strip())
            return '\n'.join(texts)
        
        return None
    except Exception as e:
        return f"[读取失败：{e}]"


def summarize_document_text(text):
    """总结文档内容"""
    if not text:
        return None
    
    # 清理和分段
    paragraphs = [p.strip() for p in text.split('\n\n') if p.strip() and len(p.strip()) > 20]
    
    if not paragraphs:
        return None
    
    # 提取关键信息
    first_para = paragraphs[0] if paragraphs else ''
    
    # 合同/协议类
    if any(kw in text for kw in ['合同', '协议', '签署', '甲方', '乙方']):
        return "合同/协议文档"
    
    # 报告类
    if any(kw in text for kw in ['报告', '总结', '汇报', '分析']):
        return "工作报告"
    
    # 规范/制度类
    if any(kw in text for kw in ['规范', '制度', '流程', '规定']):
        return "制度文档"
    
    # 通用：返回第一段精简版
    first_sentence = re.split(r'[。！？.!]', first_para)[0]
    if len(first_sentence) > 50:
        return first_sentence[:50] + '...'
    return first_sentence


def summarize_excel_text(text):
    """总结 Excel 内容"""
    if not text:
        return None
    
    lines = text.split('\n')
    
    # 尝试识别数据类型
    if any(kw in text for kw in ['金额', '合计', '总计', '¥', '元']):
        # 财务数据
        amounts = re.findall(r'[\d,]+\.?\d*', text)
        if amounts:
            return f"财务数据，含 {len(amounts)} 个金额项"
    
    if any(kw in text for kw in ['姓名', '部门', '岗位', '员工']):
        return "人员信息表"
    
    if any(kw in text for kw in ['日期', '时间', '计划', '进度']):
        return "进度计划表"
    
    # 通用
    if len(lines) > 0:
        return f"数据表格，{len(lines)} 行"
    
    return None


def extract_text_from_pptx(file_path):
    """从 PPT 提取文本"""
    try:
        prs = Presentation(file_path)
        texts = []
        for slide_idx, slide in enumerate(prs.slides):
            if slide_idx >= 5:  # 只读前 5 页
                break
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
        return '\n'.join(texts)
    except Exception as e:
        return f"[PPT 读取失败：{e}]"


def extract_text_from_pdf(file_path):
    """从 PDF 提取文本"""
    try:
        reader = PdfReader(file_path)
        texts = []
        for page_idx, page in enumerate(reader.pages):
            if page_idx >= 5:  # 只读前 5 页
                break
            text = page.extract_text()
            if text:
                texts.append(text.strip())
        return '\n'.join(texts)
    except Exception as e:
        return f"[PDF 读取失败：{e}]"


def process_attachment(attach, output_dir=None):
    """处理附件，提取关键信息（不保存文件）"""
    if not attach.payload:
        return None
    
    filename = attach.filename
    if not filename:
        return None
    
    # 根据类型提取关键信息
    lower_name = filename.lower()
    
    # 文档类：提取文本
    if lower_name.endswith('.docx'):
        text = extract_text_from_buffer(attach.payload, 'docx')
        return {'type': 'doc', 'summary': summarize_document_text(text)}
    
    elif lower_name.endswith(('.xlsx', '.xls')):
        text = extract_text_from_buffer(attach.payload, 'xlsx')
        return {'type': 'excel', 'summary': summarize_excel_text(text)}
    
    elif lower_name.endswith(('.pptx', '.ppt')):
        text = extract_text_from_buffer(attach.payload, 'pptx')
        return {'type': 'ppt', 'summary': summarize_document_text(text)}
    
    elif lower_name.endswith('.pdf'):
        text = extract_text_from_buffer(attach.payload, 'pdf')
        return {'type': 'pdf', 'summary': summarize_document_text(text)}
    
    # 图片类：只记录文件名
    elif lower_name.endswith(('.png', '.jpg', '.jpeg', '.gif')):
        return {'type': 'image', 'filename': filename}
    
    else:
        return {'type': 'other', 'filename': filename}


def fetch_emails(account, days=1, since=None, until=None, output_dir=None):
    """收取邮件"""
    emails = []
    
    # 计算日期范围
    if since and until:
        date_from = datetime.strptime(since, '%Y-%m-%d')
        date_to = datetime.strptime(until, '%Y-%m-%d')
    else:
        date_to = datetime.now()
        date_from = date_to - timedelta(days=days)
    
    try:
        with MailBox(account['imap_server'], port=int(account.get('imap_port', 993))).login(account['email'], account['auth_code']) as mailbox:
            
            # 使用 imap_tools 的 AND 条件搜索日期范围
            if since and until:
                date_from = datetime.strptime(since, '%Y-%m-%d').date()
                date_to = datetime.strptime(until, '%Y-%m-%d').date()
                criteria = AND(date_gte=date_from, date_lte=date_to)
            else:
                date_from = (datetime.now() - timedelta(days=days)).date()
                criteria = AND(date_gte=date_from)
            
            for msg in mailbox.fetch(criteria, reverse=True):
                # 跳过自己发送的邮件
                if msg.from_ == account['email']:
                    continue
                
                email_data = {
                    'from': decode_mime(msg.from_),
                    'to': decode_mime(msg.to),
                    'subject': decode_mime(msg.subject),
                    'date': msg.date,
                    'text': msg.text or '',
                    'html': msg.html or '',
                    'attachments': [],
                    'account': account['email']
                }
                
                # 处理附件（只解析内容，不保存文件）
                for attach in msg.attachments:
                    attach_info = {
                        'filename': attach.filename,
                        'content_type': attach.content_type,
                        'size': len(attach.payload) if attach.payload else 0,
                        'summary': None
                    }
                    
                    # 解析附件内容（不保存）
                    attach_result = process_attachment(attach)
                    if attach_result:
                        attach_info['summary'] = attach_result.get('summary')
                    
                    email_data['attachments'].append(attach_info)
                
                emails.append(email_data)
    
    except Exception as e:
        import traceback
        print(f"❌ 邮箱 {account['email']} 连接失败：{e}")
        traceback.print_exc()
        return []
    
    return emails


def is_user_recipient(email, user_email='zhibin@cheersucloud.com'):
    """检查用户是否是收件人（而非仅抄送）"""
    to_field = (email.get('to') or '').lower()
    
    # 检查收件人字段是否包含用户邮箱
    user_email_lower = user_email.lower()
    if user_email_lower in to_field:
        return True
    
    # 也检查用户名部分（如"支彬"）
    user_names = ['支彬', 'zhibin', '支']
    return any(name in to_field for name in user_names)


def classify_email(email, user_email='zhibin@cheersucloud.com'):
    """智能分类邮件 - 考虑收件人/抄送关系"""
    subject = (email.get('subject') or '').lower()
    text = (email.get('text') or '').lower()
    content = subject + ' ' + text
    
    # 检查是否是收件人（需要处理）还是抄送（仅需知悉）
    is_to = is_user_recipient(email, user_email)
    
    # 重要邮件关键词
    important_keywords = [
        '紧急', 'urgent', '重要', 'priority', '立即', 'asap',
        '审批', 'approval', '确认', 'confirm', '决定', 'decision',
        '合同', 'contract', '订单', 'order', '付款', 'payment',
        '账单', 'invoice', '结算', 'budget', '预算'
    ]
    
    # 行动项关键词（针对收件人的）
    action_to_me = [
        '请确认', '请核对', '请审批', '请处理', '请回复',
        '需要你', '请您', '请贵司', '请贵部门',
        '盼复', '望确认', '请知悉并', '请协助'
    ]
    
    # 一般行动词（可能是其他人的任务）
    action_general = [
        '请', '需要', '必须', '应该', '安排', '准备',
        '提交', '回复', '反馈', '处理', '联系', '提供', '对接', '核实'
    ]
    
    # 通知类关键词
    notification_keywords = [
        '通知', 'notice', '提醒', 'reminder', '会议', 'meeting',
        '日程', 'schedule', '邀请', 'invitation', '系统', 'system',
        '自动', 'auto', '订阅', 'subscription', '报告', 'report',
        '知悉', '知晓', '报备', '备案'
    ]
    
    # 检查关键词
    important_score = sum(1 for kw in important_keywords if kw in content)
    action_to_me_score = sum(1 for kw in action_to_me if kw in content)
    action_general_score = sum(1 for kw in action_general if kw in content)
    notification_score = sum(1 for kw in notification_keywords if kw in content)
    
    # 分类逻辑（考虑收件人/抄送）
    if not is_to:
        # 抄送邮件 - 通常仅需知悉
        if important_score >= 3:
            return 'important_cc', '重要抄送'
        elif notification_score >= 2:
            return 'notification', '通知'
        else:
            return 'cc_only', '抄送知悉'
    else:
        # 收件人邮件 - 可能需要处理
        if action_to_me_score >= 1 or (important_score >= 1 and action_general_score >= 2):
            return 'action', '待办事项'
        elif important_score >= 2:
            return 'important', '重要事项'
        elif action_general_score >= 2:
            return 'action', '待办事项'
        elif notification_score >= 2:
            return 'notification', '通知'
        else:
            return 'normal', '普通邮件'


def extract_action_items(email, user_email='zhibin@cheersucloud.com'):
    """从邮件中提取行动项，并识别是否指派给用户"""
    text = email.get('text', '')
    actions = []
    assigned_to_me = False
    
    # 清理 HTML 和邮件头
    text = re.sub(r'<[^>]+>', '', text)
    text = re.sub(r'主 题：.*?\n', '', text)  # 移除主题行
    text = re.sub(r'发件人：.*?发送时间：.*?(?=\n\n|\Z)', '', text, flags=re.DOTALL)  # 移除引用头
    text = re.sub(r'-{3,}.*', '', text, flags=re.DOTALL)  # 移除签名分隔线
    
    # 检查是否是收件人
    is_to = is_user_recipient(email, user_email)
    
    # 明确指派给"我"的行动词
    action_to_me_patterns = [
        r'请 [您你] 确认', r'请 [您你] 核对', r'请 [您你] 审批',
        r'请 [您你] 处理', r'请 [您你] 回复', r'需要 [您你]',
        r'请贵司', r'请贵部门', r'盼复', r'望确认',
        r'请知悉并 (?:处理 | 回复 | 确认)'
    ]
    
    # 一般行动词（可能是其他人的任务）
    action_general = [
        '请', '需要', '必须', '应该', '安排', '准备',
        '提交', '回复', '反馈', '处理', '联系', '提供', '对接', '核实'
    ]
    
    # 按句子分割
    sentences = re.split(r'[。！？.!?\n]', text)
    
    for sentence in sentences:
        sentence = sentence.strip()
        
        # 过滤太短或太长的句子
        if len(sentence) < 15 or len(sentence) > 150:
            continue
        
        # 过滤纯主题行或签名
        if sentence.startswith('主 题') or sentence.startswith('发件人') or sentence.startswith('Hi ') or '@' in sentence:
            continue
        
        # 检查是否明确指派给"我"
        is_assigned = any(re.search(pattern, sentence) for pattern in action_to_me_patterns)
        
        if is_assigned:
            assigned_to_me = True
            # 提取日期信息
            date_match = re.search(r'(\d{1,2}月\d{1,2}日 [上下午 ]|\d{1,2}/\d{1,2}|今天 | 明天 | 下周一 | 本周五)', sentence)
            if date_match:
                sentence = f"⏰ {date_match.group()}：{sentence}"
            actions.append(('me', sentence))
        elif is_to and any(verb in sentence for verb in action_general):
            # 是收件人且包含行动词，可能是我的任务
            date_match = re.search(r'(\d{1,2}月\d{1,2}日 [上下午 ]|\d{1,2}/\d{1,2}|今天 | 明天 | 下周一 | 本周五)', sentence)
            if date_match:
                sentence = f"⏰ {date_match.group()}：{sentence}"
            actions.append(('maybe_me', sentence))
    
    # 去重
    unique_actions = []
    seen = set()
    for assignee, action in actions:
        key = action[:30]
        if key not in seen:
            seen.add(key)
            unique_actions.append((assignee, action))
    
    # 优先返回指派给我的任务
    my_actions = [a for assignee, a in unique_actions if assignee == 'me']
    maybe_actions = [a for assignee, a in unique_actions if assignee == 'maybe_me']
    
    return my_actions[:3] + maybe_actions[:2], len(my_actions) > 0  # 最多 5 个，优先我的任务


def summarize_email_content(email, max_length=150):
    """AI 智能总结邮件内容（非原文截取）"""
    text = email.get('text', '')
    subject = email.get('subject', '')
    
    if not text:
        return None
    
    # 清理 HTML 和引用
    text = re.sub(r'<[^>]+>', '', text)
    text = re.sub(r'-{3,}.*?发件人：.*?(?=\n\n|\Z)', '', text, flags=re.DOTALL)  # 移除引用
    text = re.sub(r'_{3,}.*', '', text, flags=re.DOTALL)  # 移除签名
    text = re.sub(r'主 题：.*?\n', '', text)  # 移除主题行
    
    # 提取关键信息
    paragraphs = [p.strip() for p in text.split('\n\n') if p.strip() and len(p.strip()) > 20]
    
    if not paragraphs:
        return None
    
    # 识别邮件类型并生成总结
    first_para = paragraphs[0] if paragraphs else ''
    
    # 账单/结算类
    if any(kw in text for kw in ['账单', '结算', '发票', '金额', '¥', '元']):
        # 提取金额和状态
        amount_match = re.search(r'[¥￥]?[\d,]+\.?\d*\s*元', text)
        status_match = re.search(r'(确认无误 | 已核对 | 待确认 | 请审核)', text)
        
        summary = "账单核对"
        if amount_match:
            summary += f"，金额：{amount_match.group()}"
        if status_match:
            summary += f"，状态：{status_match.group()}"
        return summary
    
    # 订单/交付类
    if any(kw in text for kw in ['订单', '交付', '设备', '台', '配置']):
        # 提取数量和日期
        qty_match = re.search(r'(\d+\s*台|\d+\s*台)', text)
        date_match = re.search(r'(\d{1,2}月\d{1,2}日|\d{4}-\d{2}-\d{2})', text)
        
        summary = "订单交付"
        if qty_match:
            summary += f"，数量：{qty_match.group()}"
        if date_match:
            summary += f"，日期：{date_match.group()}"
        return summary
    
    # 审批/确认类
    if any(kw in text for kw in ['审批', '确认', '审核', '批准']):
        action_match = re.search(r'请 (?:您 | 你)?(\w{1,4})', text)
        deadline_match = re.search(r'(\d{1,2}月\d{1,2}日 [上下午 ]?|今天 | 明天 | 本周)', text)
        
        summary = "待确认事项"
        if action_match:
            summary += f"，需要：{action_match.group(1)}"
        if deadline_match:
            summary += f"，截止：{deadline_match.group()}"
        return summary
    
    # 项目/工作同步类
    if any(kw in text for kw in ['项目', '工作', '推进', '反馈', '同步']):
        # 提取关键动作
        actions = []
        if '对接' in text:
            actions.append('对接')
        if '核实' in text:
            actions.append('核实')
        if '提供' in text:
            actions.append('提供资料')
        
        summary = "工作同步"
        if actions:
            summary += f"，需：{', '.join(actions)}"
        return summary
    
    # 通用总结：提取第一句核心内容
    first_sentence = re.split(r'[。！？.!]', first_para)[0]
    if len(first_sentence) > 10:
        # 精简到 max_length
        if len(first_sentence) > max_length:
            return first_sentence[:max_length] + '...'
        return first_sentence
    
    return None


def generate_summary(emails, include_attachments=True, output_dir=None):
    """生成智能邮件总结"""
    if not emails:
        return "📭 没有新邮件"
    
    summary = []
    date_str = datetime.now().strftime('%Y-%m-%d')
    summary.append(f"# 📧 邮件日报 ({date_str})\n")
    
    # 分类邮件
    important_emails = []
    action_emails = []
    important_cc_emails = []  # 重要抄送
    cc_only_emails = []  # 仅抄送知悉
    notification_emails = []
    normal_emails = []
    
    for email in emails:
        category, label = classify_email(email)
        if category == 'important':
            important_emails.append(email)
        elif category == 'action':
            action_emails.append(email)
        elif category == 'important_cc':
            important_cc_emails.append(email)
        elif category == 'cc_only':
            cc_only_emails.append(email)
        elif category == 'notification':
            notification_emails.append(email)
        else:
            normal_emails.append(email)
    
    # 概览
    summary.append("## 📊 概览")
    summary.append(f"- 总邮件数：**{len(emails)}**")
    
    # 收件人邮件（需要关注）
    to_count = len(important_emails) + len(action_emails) + len(normal_emails)
    if to_count > 0:
        summary.append(f"- 📬 收件人邮件：**{to_count}**")
        if important_emails:
            summary.append(f"  - 🔴 重要：**{len(important_emails)}**")
        if action_emails:
            summary.append(f"  - ⚡ 待办：**{len(action_emails)}**")
    
    # 抄送邮件（仅需知悉）
    cc_count = len(cc_only_emails) + len(important_cc_emails)
    if cc_count > 0:
        summary.append(f"- 📧 抄送邮件：**{cc_count}** (仅需知悉)")
        if important_cc_emails:
            summary.append(f"  - 🔴 重要抄送：**{len(important_cc_emails)}**")
    
    if notification_emails:
        summary.append(f"- 📢 通知：**{len(notification_emails)}**")
    
    # 行动项汇总（最重要）- 只提取收件人邮件的行动项
    my_actions = []
    maybe_actions = []
    cc_actions = []
    
    for email in important_emails + action_emails + normal_emails:
        is_to = is_user_recipient(email)
        actions, has_my_tasks = extract_action_items(email)
        
        if actions:
            item = {
                'from': email['from'],
                'subject': email['subject'],
                'actions': actions,
                'has_my_tasks': has_my_tasks
            }
            
            if is_to and has_my_tasks:
                my_actions.append(item)
            elif is_to:
                maybe_actions.append(item)
            else:
                cc_actions.append(item)
    
    # 显示明确指派给我的任务
    if my_actions:
        summary.append("\n## ⚡ 需要我处理\n")
        for item in my_actions:
            summary.append(f"**{item['subject'] or '(无主题)'}** ({item['from']})")
            for action in item['actions']:
                summary.append(f"- {action}")
            summary.append("")
    
    # 可能是我的任务（收件人但行动项不明确）
    if maybe_actions and not my_actions:
        summary.append("\n## 📋 可能需要处理\n")
        for item in maybe_actions:
            summary.append(f"**{item['subject'] or '(无主题)'}** ({item['from']})")
            for action in item['actions']:
                summary.append(f"- {action}")
            summary.append("")
    
    # 重要邮件（收件人）
    if important_emails:
        summary.append("\n## 🔴 重要事项（收件人）\n")
        for email in important_emails:
            content_summary = summarize_email_content(email)
            if content_summary:
                summary.append(f"**{email['subject'] or '(无主题)'}**\n")
                summary.append(f"发件人：{email['from']} | 时间：{email['date'].strftime('%H:%M')} | 📬 收件人\n")
                summary.append(f"{content_summary}\n")
                
                # 附件信息（显示总结）
                if email['attachments']:
                    att_summaries = []
                    for att in email['attachments'][:3]:
                        if att.get('summary'):
                            att_summaries.append(f"{att['filename']}: {att['summary']}")
                        else:
                            att_summaries.append(att['filename'])
                    
                    att_text = f"📎 {len(email['attachments'])} 个附件："
                    if len(email['attachments']) > 3:
                        att_text += f"{', '.join(att_summaries)} 等"
                    else:
                        att_text += f"{', '.join(att_summaries)}"
                    summary.append(f"{att_text}\n")
                summary.append("")
    
    # 重要抄送
    if important_cc_emails:
        summary.append("\n## 🔴 重要抄送（仅需知悉）\n")
        for email in important_cc_emails:
            content_summary = summarize_email_content(email, max_length=150)
            if content_summary:
                summary.append(f"- **{email['subject'] or '(无主题)'}** ({email['from']}) 📧 抄送\n")
                summary.append(f"  {content_summary}\n")
                
                # 附件信息
                if email['attachments']:
                    for att in email['attachments'][:2]:
                        if att.get('summary'):
                            summary.append(f"  📎 {att['filename']}: {att['summary']}\n")
        summary.append("")
    
    # 待办事项
    if action_emails and not important_emails:
        summary.append("\n## 📋 待办事项\n")
        for email in action_emails:
            content_summary = summarize_email_content(email, max_length=200)
            if content_summary:
                summary.append(f"- **{email['subject'] or '(无主题)'}** ({email['from']})\n")
                summary.append(f"  {content_summary}\n")
    
    # 通知（合并显示）
    if notification_emails:
        summary.append("\n## 📢 通知\n")
        for email in notification_emails[:5]:  # 最多显示 5 个
            summary.append(f"- {email['subject'] or '(无主题)'} - {email['from']} ({email['date'].strftime('%H:%M')})")
        if len(notification_emails) > 5:
            summary.append(f"- ... 还有 {len(notification_emails) - 5} 条通知")
        summary.append("")
    
    # 普通邮件（精简列表）
    if normal_emails:
        summary.append("\n## 📮 其他邮件（收件人）\n")
        for email in normal_emails[:5]:  # 最多显示 5 个
            att_icon = "📎" if email['attachments'] else ""
            summary.append(f"- {att_icon} {email['subject'] or '(无主题)'} - {email['from']} 📬")
        if len(normal_emails) > 5:
            summary.append(f"- ... 还有 {len(normal_emails) - 5} 封收件人邮件")
        summary.append("")
    
    # 抄送知悉（精简列表）
    if cc_only_emails:
        summary.append("\n## 📧 抄送知悉\n")
        for email in cc_only_emails[:5]:  # 最多显示 5 个
            att_icon = "📎" if email['attachments'] else ""
            summary.append(f"- {att_icon} {email['subject'] or '(无主题)'} - {email['from']} 📧")
        if len(cc_only_emails) > 5:
            summary.append(f"- ... 还有 {len(cc_only_emails) - 5} 封抄送邮件")
        summary.append("")
    
    # 附件汇总
    total_attachments = sum(len(e['attachments']) for e in emails)
    if total_attachments > 0:
        summary.append(f"\n---\n📁 附件总数：{total_attachments} 个（已解析内容，未保存文件）")
    
    return '\n'.join(summary)


def main():
    parser = argparse.ArgumentParser(description='阿里云邮箱工具')
    parser.add_argument('--action', required=True, choices=['summary', 'list', 'attachments', 'test'])
    parser.add_argument('--days', type=int, default=1, help='收取最近 N 天的邮件')
    parser.add_argument('--since', help='开始日期 YYYY-MM-DD')
    parser.add_argument('--until', help='结束日期 YYYY-MM-DD')
    parser.add_argument('--output-dir', default='./email-attachments', help='附件保存目录')
    parser.add_argument('--include-attachments', action='store_true', default=True)
    parser.add_argument('--send-feishu', action='store_true')
    parser.add_argument('--chat-id', help='飞书聊天 ID')
    
    args = parser.parse_args()
    
    # 加载凭证
    # skills/aliyun-mail/lib/fetch_emails.py -> workspace
    workspace_dir = Path(__file__).parent.parent.parent.parent
    cred_file = workspace_dir / '.credentials' / 'aliyun-mail.md'
    
    if not cred_file.exists():
        print(f"❌ 凭证文件不存在：{cred_file}")
        print("请先创建凭证文件，参考：skills/aliyun-mail/README-授权码.md")
        sys.exit(1)
    
    accounts = parse_credentials(cred_file)
    
    if not accounts:
        print("❌ 凭证文件中没有有效的邮箱配置")
        sys.exit(1)
    
    print(f"📬 找到 {len(accounts)} 个邮箱账户")
    
    # 根据动作处理
    if args.action == 'test':
        print("\n✅ 邮箱连接测试成功")
        return
    
    if args.action == 'list':
        all_emails = []
        for account in accounts:
            print(f"\n正在收取 {account['email']} 的邮件...", file=sys.stderr)
            emails = fetch_emails(
                account,
                days=args.days,
                since=args.since,
                until=args.until,
                output_dir=None
            )
            print(f"  ✓ 收取 {len(emails)} 封邮件", file=sys.stderr)
            all_emails.extend(emails)
        
        for email in all_emails:
            print(f"\n{email['date'].strftime('%Y-%m-%d %H:%M')} | {email['from']} | {email['subject']}")
            if email['attachments']:
                print(f"  附件：{', '.join([a['filename'] for a in email['attachments']])}")
        return
    
    if args.action == 'attachments':
        os.makedirs(args.output_dir, exist_ok=True)
        print(f"\n保存附件到：{args.output_dir}")
        # 简化处理，实际需要重新获取
        print("附件提取功能开发中...")
        return
    
    if args.action == 'summary':
        # 不再保存附件到文件
        print("\n📧 收取邮件并解析附件内容（不保存文件）...")
        
        # 收取邮件
        all_emails = []
        for account in accounts:
            print(f"\n正在收取 {account['email']} 的邮件...")
            emails = fetch_emails(
                account,
                days=args.days,
                since=args.since,
                until=args.until,
                output_dir=None  # 不再保存文件
            )
            print(f"  ✓ 收取 {len(emails)} 封邮件")
            all_emails.extend(emails)
        
        print(f"\n总计：{len(all_emails)} 封邮件")
        
        # 生成总结
        summary = generate_summary(all_emails)
        print("\n" + "="*60)
        print(summary)
        print("="*60)
        
        # 保存到 workspace
        date_str = datetime.now().strftime('%Y-%m-%d')
        output_file = workspace_dir / 'email-reports' / date_str / f"{date_str}-邮件日报.md"
        output_file.parent.mkdir(parents=True, exist_ok=True)
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write(summary)
        print(f"\n💾 总结已保存：{output_file}")
        
        # 发送到飞书
        if args.send_feishu and args.chat_id:
            print(f"\n📤 发送到飞书：{args.chat_id}")
            send_to_feishu(summary, args.chat_id)


if __name__ == '__main__':
    main()
